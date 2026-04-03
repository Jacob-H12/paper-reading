#!/usr/bin/env python3
"""
AD 脑控游戏项目 — 调研汇报 PPT 生成脚本
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============ 颜色方案 ============
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MED = RGBColor(0x16, 0x21, 0x3E)
BG_LIGHT = RGBColor(0x1F, 0x2B, 0x4D)
ACCENT_BLUE = RGBColor(0x00, 0xD2, 0xFF)
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43)
ACCENT_RED = RGBColor(0xFF, 0x6B, 0x6B)
ACCENT_PURPLE = RGBColor(0xA2, 0x9B, 0xFE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xBB)
LIGHT_GRAY = RGBColor(0xDD, 0xDD, 0xEE)


def add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha:
        shape.fill.fore_color.brightness = alpha
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, items, left, top, width, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = spacing
    return txBox


# ============ Slide 1: 封面 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# 装饰线
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(13.333), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "面向阿尔茨海默患者的\nEEG 闭环游戏干预系统", 42, ACCENT_BLUE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
             "— 系统调研汇报 —", 24, GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
             "调研方向：AD 稳定 EEG 指标 · 游戏干预框架 · 闭环控制架构", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.5),
             datetime.datetime.now().strftime("%Y年%m月%d日"), 16, GRAY, False, PP_ALIGN.CENTER)


# ============ Slide 2: 目录 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7), "汇报大纲", 32, ACCENT_BLUE, True)

sections = [
    ("01", "研究背景与动机", ACCENT_BLUE),
    ("02", "AD 患者稳定 EEG 指标", ACCENT_GREEN),
    ("03", "关键发现：Alpha 双相轨迹", ACCENT_ORANGE),
    ("04", "文献中的游戏干预方案", ACCENT_PURPLE),
    ("05", "游戏设计框架建议", ACCENT_BLUE),
    ("06", "闭环控制架构", ACCENT_GREEN),
    ("07", "下一步计划", ACCENT_ORANGE),
]

for i, (num, title, color) in enumerate(sections):
    y = Inches(1.5) + Inches(0.75) * i
    add_shape_bg(slide, Inches(1.5), y, Inches(10), Inches(0.6), BG_LIGHT)
    add_text_box(slide, Inches(1.7), y + Pt(4), Inches(1), Inches(0.5), num, 22, color, True)
    add_text_box(slide, Inches(2.7), y + Pt(4), Inches(8), Inches(0.5), title, 20, WHITE, False)


# ============ Slide 3: 研究背景 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "01 研究背景与动机", 32, ACCENT_BLUE, True)

# 左列
add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), BG_LIGHT)
add_text_box(slide, Inches(1.0), Inches(1.6), Inches(5.4), Inches(0.5), "阿尔茨海默病 (AD)", 22, ACCENT_BLUE, True)
add_bullet_slide_content(slide, [
    "• 全球预计 2050 年影响 1.315 亿人",
    "• 经济负担：全球 $9.1 万亿/年",
    "• 病理机制：Aβ 沉积 → tau 过磷酸化 → 突触丧失 → 神经退行",
    "• 临床前期病变可提前 20 年开始",
    "• 目前无法治愈，干预越早效果越好",
], Inches(1.0), Inches(2.3), Inches(5.4), 16)

# 右列
add_shape_bg(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), BG_LIGHT)
add_text_box(slide, Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.5), "为什么用 EEG + 游戏？", 22, ACCENT_GREEN, True)
add_bullet_slide_content(slide, [
    "• EEG：非侵入、低成本、毫秒级时间分辨率",
    "• 专家共识推荐 rsEEG 用于 AD 临床试验 (Babiloni 2021)",
    "• 游戏化：提高患者依从性和参与度",
    "• 闭环：实时神经反馈，个性化干预",
    "• 文献中已有 5+ 种闭环游戏报告有效",
], Inches(7.1), Inches(2.3), Inches(5.4), 16)


# ============ Slide 4: 核心文献 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7), "核心参考文献", 32, ACCENT_BLUE, True)

# 文献1
add_shape_bg(slide, Inches(0.8), Inches(1.4), Inches(5.8), Inches(2.5), BG_LIGHT)
add_text_box(slide, Inches(1.0), Inches(1.5), Inches(5.4), Inches(0.4), "📄 文献 1 — 专家共识 (Babiloni et al. 2021)", 18, ACCENT_BLUE, True)
add_bullet_slide_content(slide, [
    "Alzheimers & Dementia, IF≈14, 60+ 专家联合署名",
    "推荐 rsEEG 指标用于 AD 临床试验",
    "核心：α 功率/峰频 + δ/θ 功率 + 连接性(DTF/PLI/LLC)",
    "标准化记录与分析流程",
], Inches(1.0), Inches(2.1), Inches(5.4), 14, LIGHT_GRAY)

# 文献2
add_shape_bg(slide, Inches(6.9), Inches(1.4), Inches(5.8), Inches(2.5), BG_LIGHT)
add_text_box(slide, Inches(7.1), Inches(1.5), Inches(5.4), Inches(0.4), "📄 文献 2 — 系统综述 (Sohrabpour et al. 2025)", 18, ACCENT_GREEN, True)
add_bullet_slide_content(slide, [
    "Alzheimers & Dementia, 49 篇多模态研究系统综述",
    "EEG/MEG + PET/CSF 蛋白标志物联合分析",
    "核心发现：α 振荡双相轨迹（先升后降）",
    "频段-病理特异性：δ/θ↔Aβ, α/β↔Aβ+tau",
], Inches(7.1), Inches(2.1), Inches(5.4), 14, LIGHT_GRAY)

# 其他文献
add_shape_bg(slide, Inches(0.8), Inches(4.2), Inches(11.9), Inches(2.8), BG_LIGHT)
add_text_box(slide, Inches(1.0), Inches(4.3), Inches(11.5), Inches(0.4), "📚 闭环游戏相关文献 (7篇)", 18, ACCENT_ORANGE, True)
add_bullet_slide_content(slide, [
    "[Cancer 2024] Prinsloo — BCI 游戏治疗 CIPN，双盲 RCT，d=1.07",
    "[Front Hum Neurosci 2019] Arvaneh — P300 自适应 BCI 注意力训练，30min 即见效",
    "[J Healthc Eng 2021] Wang — 3 种严肃游戏 (拔河/小鸟/果冻)，5 级注意力闭环",
    "[J Neuroeng Rehabil 2026] Yan — 平板 Go/NoGo + 单通道 EEG 神经反馈 (ADHD)",
    "[Neuroscience 2026] Yang — Unity SkiSport 滑雪游戏 + XGBoost 自适应难度 (ADHD)",
], Inches(1.0), Inches(4.9), Inches(11.5), 14, LIGHT_GRAY, Pt(4))


# ============ Slide 5: 确定指标总览 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "02 AD 患者稳定 EEG 指标 — 总览", 32, ACCENT_GREEN, True)

# 表格
headers = ["#", "指标", "状态", "共识级别", "效应量(d)", "实时性", "推荐"]
rows = [
    ["1", "频谱慢化 (α↓+δ/θ↑)", "静息", "⭐⭐⭐⭐⭐", "0.8-1.5", "⚡即时", "✅ 确定"],
    ["2", "SPR = (α+β)/(δ+θ)", "静息+任务", "⭐⭐⭐⭐", "1.0-1.8", "⚡即时", "✅ 确定"],
    ["3", "Alpha 峰频率↓ (APF)", "静息", "⭐⭐⭐⭐⭐", "0.9-1.3", "⚡即时", "✅ 确定"],
    ["4", "P300 潜伏期↑ 振幅↓", "任务", "⭐⭐⭐⭐⭐", "1.0-1.5", "⏱1-3s", "✅ 确定"],
    ["5", "Beta 功率↓", "静息", "⭐⭐⭐⭐", "0.7-1.0", "⚡即时", "✅ 确定"],
    ["6", "α/β 功能连接↓", "静息", "⭐⭐⭐⭐", "0.6-1.0", "⏱2-4s", "🔶 大概率"],
    ["7", "Delta 功率↑", "静息", "⭐⭐⭐⭐", "0.5-1.0", "⚡即时", "🔶 大概率"],
    ["8", "微状态 C↓", "静息", "⭐⭐⭐", "—", "⏱阶段性", "🔶 大概率"],
]

col_widths = [Inches(0.4), Inches(3.0), Inches(1.2), Inches(1.5), Inches(1.3), Inches(1.0), Inches(1.2)]
table_left = Inches(0.8)
table_top = Inches(1.4)
row_h = Inches(0.55)

# header
for j, (hdr, w) in enumerate(zip(headers, col_widths)):
    x = table_left + sum(col_widths[:j], Emu(0))
    shape = add_shape_bg(slide, x, table_top, w - Pt(2), row_h, ACCENT_BLUE)
    add_text_box(slide, x + Pt(4), table_top + Pt(4), w - Pt(8), row_h, hdr, 13, BG_DARK, True, PP_ALIGN.CENTER)

# rows
for i, row in enumerate(rows):
    y = table_top + row_h * (i + 1) + Pt(2)
    bg_color = BG_LIGHT if i % 2 == 0 else BG_MED
    for j, (cell, w) in enumerate(zip(row, col_widths)):
        x = table_left + sum(col_widths[:j], Emu(0))
        add_shape_bg(slide, x, y, w - Pt(2), row_h - Pt(2), bg_color)
        c = ACCENT_GREEN if "✅" in cell else (ACCENT_ORANGE if "🔶" in cell else WHITE)
        add_text_box(slide, x + Pt(4), y + Pt(4), w - Pt(8), row_h, cell, 12, c, False, PP_ALIGN.CENTER)


# ============ Slide 6: 频谱慢化详解 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "指标 1: 频谱慢化 — 最经典的 AD EEG 标志", 28, ACCENT_GREEN, True)

add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5), BG_LIGHT)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5.4), Inches(0.4), "现象与机制", 20, ACCENT_BLUE, True)
add_bullet_slide_content(slide, [
    "• 后部 Alpha (8-12Hz) 功率 显著降低",
    "• Delta (<4Hz) / Theta (4-8Hz) 功率 广泛增加",
    "• 整体频谱向慢波方向偏移",
    "",
    "机制 (Babiloni 2021 生理模型):",
    "  正常：皮层-丘脑-基底前脑-脑干网络",
    "       → 主导 Alpha 节律 → 调节唤醒/意识",
    "  AD：胆碱能退化 → ACh↓ → 低频活动↑",
    "       丘脑-皮层 \"断连模式\"",
    "       → Alpha 去同步化 + Delta/Theta 异常增强",
], Inches(1.0), Inches(2.0), Inches(5.4), 14, LIGHT_GRAY)

add_shape_bg(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.5), BG_LIGHT)
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.4), Inches(0.4), "数据支撑", 20, ACCENT_GREEN, True)
add_bullet_slide_content(slide, [
    "文献支持：40+ 年持续复现",
    "  - Sohrabpour 2025: 10+ 篇独立报告 δ/θ↑",
    "  - Sohrabpour 2025: 8+ 篇独立报告 α↓",
    "  - Babiloni 2021: 专家共识一级推荐",
    "",
    "效应量：",
    "  - Alpha 下降: Cohen's d = 0.8-1.5",
    "  - Delta 增加: Cohen's d = 0.8-1.2 (中晚期)",
    "",
    "认知关联：",
    "  - δ/θ 功率 ↔ MMSE 负相关",
    "  - α 功率 ↔ MMSE 正相关",
    "",
    "关键通道：P3, P4, O1, O2, Pz (后部)",
], Inches(7.1), Inches(2.0), Inches(5.4), 14, LIGHT_GRAY)


# ============ Slide 7: SPR + APF ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "指标 2-3: SPR 综合比值 + Alpha 峰频率", 28, ACCENT_GREEN, True)

# SPR
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5), BG_LIGHT)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5.4), Inches(0.4), "SPR = (α+β) / (δ+θ)", 22, ACCENT_BLUE, True)
add_bullet_slide_content(slide, [
    "• 整合四个频段的综合指标",
    "• 效应量 d = 1.0-1.8（比单频段更鲁棒）",
    "• Aβ+ aMCI 中显著降低",
    "",
    "⭐ 核心优势：",
    "  避免了 Alpha 双相轨迹的陷阱",
    "  即使 α 代偿性↑，若 δ/θ 也在↑",
    "  SPR 仍会下降 → 更可靠",
    "",
    "POSI (病理性振荡慢化指数)：",
    "  考虑年龄匹配健康对照的偏移",
    "  与认知评分强负相关",
    "",
    "→ 推荐作为闭环系统的主控变量",
], Inches(1.0), Inches(2.0), Inches(5.4), 14, LIGHT_GRAY)

# APF
add_shape_bg(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.5), BG_LIGHT)
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.4), Inches(0.4), "Alpha 峰频率 (APF)", 22, ACCENT_GREEN, True)
add_bullet_slide_content(slide, [
    "• 正常人：~10-10.5 Hz",
    "• MCI 患者：~9.5 Hz",
    "• AD 患者：~8-9 Hz",
    "",
    "⭐ 核心优势：",
    "  即使 α 功率代偿性升高",
    "  APF 仍然下移",
    "  → 比 α 功率更稳健的早期指标",
    "",
    "稳定性：",
    "  test-retest ICC = 0.85-0.92",
    "  是 EEG 中最稳定的个体特征",
    "",
    "分子关联：",
    "  与 PET Aβ/tau 负相关",
    "  与 CSF Aβ42 正相关",
], Inches(7.1), Inches(2.0), Inches(5.4), 14, LIGHT_GRAY)


# ============ Slide 8: P300 + Beta ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "指标 4-5: P300 + Beta 功率", 28, ACCENT_GREEN, True)

# P300
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.8), Inches(5.5), BG_LIGHT)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(5.4), Inches(0.4), "P300 潜伏期延长 + 振幅降低", 20, ACCENT_BLUE, True)
add_bullet_slide_content(slide, [
    "• 任务态指标（需 oddball 范式）",
    "• 潜伏期延长 30-80ms，振幅降低 2-5μV",
    "• 效应量：潜伏期 d=1.0-1.5, 振幅 d=0.8-1.2",
    "",
    "认知意义：",
    "  直接反映 \"注意力分配/目标检测\" 能力",
    "  与工作记忆、反应速度强相关",
    "",
    "🎮 与游戏的天然联系：",
    "  P300 = \"检测到新异刺激\" 的脑反应",
    "  游戏中嵌入目标检测任务",
    "  （找不同、捕捉特定物品）",
    "  → 直接诱发和监测 P300",
    "  → P300 改善 = 通关/升级条件",
], Inches(1.0), Inches(2.0), Inches(5.4), 14, LIGHT_GRAY)

# Beta
add_shape_bg(slide, Inches(6.9), Inches(1.3), Inches(5.8), Inches(5.5), BG_LIGHT)
add_text_box(slide, Inches(7.1), Inches(1.4), Inches(5.4), Inches(0.4), "Beta 功率降低", 20, ACCENT_GREEN, True)
add_bullet_slide_content(slide, [
    "• 静息态，枕/颞/顶区",
    "• Beta (12-30Hz) 显著降低",
    "• 效应量：d = 0.7-1.0",
    "",
    "AD 特异性（关键）：",
    "  SNAP 患者 Beta 反而增加",
    "  AD 患者 Beta 降低",
    "  → 可用于 AD vs 非 AD 鉴别",
    "",
    "分子关联：",
    "  与 Aβ 正相关/倒 U 型",
    "  与 tau 负相关",
    "",
    "🎮 游戏映射：",
    "  Beta ↔ 主动认知参与",
    "  → 作为 \"参与度\" 指标",
], Inches(7.1), Inches(2.0), Inches(5.4), 14, LIGHT_GRAY)


# ============ Slide 9: Alpha 双相轨迹（重点！）============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "03 关键发现：Alpha 双相轨迹 ⚠", 32, ACCENT_ORANGE, True)
add_text_box(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.4), "Sohrabpour et al. (2025) 的核心突破", 18, GRAY, False)

# 轨迹图示意（文字版）
add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.5), BG_LIGHT)
add_text_box(slide, Inches(1.0), Inches(1.7), Inches(11.3), Inches(0.4), "Alpha 振荡在 AD 全病程中呈倒 U 型（双相）轨迹", 20, ACCENT_ORANGE, True)

trajectory_text = (
    "α 活动\n"
    "  ↑         ╱╲\n"
    "  │       ╱    ╲\n"
    "  │     ╱        ╲\n"
    "  │───╱── 基线 ──  ╲───────────\n"
    "  │                    ╲\n"
    "  ↓                      ╲________\n"
    "      临床前(Aβ+)   MCI    轻度AD   中重度AD"
)
add_text_box(slide, Inches(2.0), Inches(2.2), Inches(9), Inches(2.0), trajectory_text, 14, ACCENT_BLUE, False, PP_ALIGN.LEFT, "Courier New")

# 解释
add_shape_bg(slide, Inches(0.8), Inches(4.3), Inches(5.8), Inches(2.8), BG_LIGHT)
add_text_box(slide, Inches(1.0), Inches(4.4), Inches(5.4), Inches(0.4), "机制解释", 18, ACCENT_BLUE, True)
add_bullet_slide_content(slide, [
    "早期：Aβ 损伤 GABA 能抑制性中间神经元",
    "  → E/I 失衡 → 皮层过度兴奋 → α 代偿↑",
    "中期：tau 病理加速 → 突触丧失",
    "  → 神经退行 → α 开始下降",
    "晚期：广泛神经元丢失 → α 严重衰减",
], Inches(1.0), Inches(5.0), Inches(5.4), 14, LIGHT_GRAY)

add_shape_bg(slide, Inches(6.9), Inches(4.3), Inches(5.8), Inches(2.8), BG_LIGHT)
add_text_box(slide, Inches(7.1), Inches(4.4), Inches(5.4), Inches(0.4), "⚠ 对闭环设计的影响", 18, ACCENT_RED, True)
add_bullet_slide_content(slide, [
    "不能简单用 \"α↓ = AD 恶化\" 做闭环",
    "MCI/早期 AD 患者 α 可能正常甚至偏高",
    "",
    "→ 解决方案：",
    "  1. 用 SPR 代替单一 α 功率",
    "  2. 用 APF（不受双相影响）",
    "  3. 分期考虑基线差异",
], Inches(7.1), Inches(5.0), Inches(5.4), 14, LIGHT_GRAY)


# ============ Slide 10: 大概率指标 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "大概率指标（辅助评估）", 28, ACCENT_ORANGE, True)

# 3列
for i, (title, color, items) in enumerate([
    ("α/β 功能连接↓", ACCENT_BLUE, [
        "AEC/Coherence/PLV",
        "d = 0.6-1.0",
        "多种算法方向一致",
        "",
        "⭐ α/β AEC↓ 出现在",
        "认知下降之前",
        "→ 早期预警能力",
        "",
        "适合：阶段性评估",
    ]),
    ("Delta 功率↑", ACCENT_GREEN, [
        "AD 特异性标志",
        "d = 0.5-1.0",
        "SNAP 中 δ 反而↓",
        "→ AD 鉴别价值",
        "",
        "与 CSF Aβ42 负相关",
        "",
        "适合：安全守护",
        "δ过多 = 嗜睡 → 暂停游戏",
    ]),
    ("微状态变化", ACCENT_PURPLE, [
        "微状态 A/B 出现率↑",
        "微状态 C/D 出现率↓",
        "",
        "微状态 C duration",
        "与 MMSE/MoCA 正相关",
        "",
        "与 fMRI 发现高度一致",
        "",
        "适合：会话前后对比评估",
    ]),
]):
    left = Inches(0.8) + Inches(4.1) * i
    add_shape_bg(slide, left, Inches(1.3), Inches(3.8), Inches(5.5), BG_LIGHT)
    add_text_box(slide, left + Inches(0.2), Inches(1.4), Inches(3.4), Inches(0.4), title, 18, color, True)
    add_bullet_slide_content(slide, items, left + Inches(0.2), Inches(2.0), Inches(3.4), 14, LIGHT_GRAY)


# ============ Slide 11: 文献中的游戏 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "04 文献中的游戏干预方案", 32, ACCENT_PURPLE, True)

games = [
    ("EEG 控制视频游戏", "Prinsloo 2024", "CIPN 疼痛", "EEG 全频段", "RCT 双盲\nd=1.07 ✅"),
    ("P300 自适应训练", "Arvaneh 2019", "健康人注意力", "P300 能量", "30min 见效\nP300↑ Alpha↓ ✅"),
    ("拔河/小鸟/果冻 ×3", "Wang 2021", "健康人注意力", "5级注意力(IRF)", "Schulte\np<0.05 ✅"),
    ("平板 Go/NoGo", "Yan 2026", "ADHD 儿童", "单通道额叶", "NFb > 纯行为\n1月干预 ✅"),
    ("SkiSport 滑雪", "Yang 2026", "ADHD", "eSense→XGBoost", "+21.5%\nUnity 开发 ✅"),
]

headers2 = ["游戏", "来源", "目标人群", "EEG 指标", "效果"]
col_w2 = [Inches(2.8), Inches(1.8), Inches(2.0), Inches(2.4), Inches(2.6)]
t_left = Inches(0.8)
t_top = Inches(1.3)
rh = Inches(1.0)

for j, (h, w) in enumerate(zip(headers2, col_w2)):
    x = t_left + sum(col_w2[:j], Emu(0))
    add_shape_bg(slide, x, t_top, w - Pt(2), Inches(0.5), ACCENT_PURPLE)
    add_text_box(slide, x + Pt(4), t_top + Pt(4), w - Pt(8), Inches(0.5), h, 14, BG_DARK, True, PP_ALIGN.CENTER)

for i, row in enumerate(games):
    y = t_top + Inches(0.55) + rh * i
    bg = BG_LIGHT if i % 2 == 0 else BG_MED
    for j, (cell, w) in enumerate(zip(row, col_w2)):
        x = t_left + sum(col_w2[:j], Emu(0))
        add_shape_bg(slide, x, y, w - Pt(2), rh - Pt(2), bg)
        add_text_box(slide, x + Pt(6), y + Pt(6), w - Pt(12), rh, cell, 12, WHITE, False, PP_ALIGN.CENTER)

# 底部注释
add_shape_bg(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.8), BG_LIGHT)
add_text_box(slide, Inches(1.0), Inches(6.45), Inches(11.3), Inches(0.7),
             "⚠ 5/5 游戏干预均报告有效，但：① 无一直接针对 AD 患者 ② 样本量均较小(10-50人) ③ Yang(2026) 存在方法学问题(信息泄漏)",
             14, ACCENT_ORANGE, False)


# ============ Slide 12: 游戏设计框架 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "05 游戏设计框架建议", 32, ACCENT_BLUE, True)

# 4个子游戏模块
modules = [
    ("🏠 记忆花园", "情景记忆", "配对翻牌 / 物品记忆", "P300 + SPR", ACCENT_BLUE),
    ("🎵 节奏之路", "注意力/处理速度", "节奏/音乐追踪", "Alpha + APF", ACCENT_GREEN),
    ("🧩 模式拼图", "执行功能/工作记忆", "模式匹配 / 序列记忆", "SPR + 连接性", ACCENT_ORANGE),
    ("🌊 放松港湾", "情绪调节/休息", "冥想 / 呼吸引导", "Alpha + Delta", ACCENT_PURPLE),
]

for i, (name, target, proto, eeg, color) in enumerate(modules):
    left = Inches(0.6) + Inches(3.15) * i
    add_shape_bg(slide, left, Inches(1.3), Inches(3.0), Inches(3.0), BG_LIGHT)
    add_text_box(slide, left + Inches(0.15), Inches(1.4), Inches(2.7), Inches(0.5), name, 18, color, True, PP_ALIGN.CENTER)
    add_bullet_slide_content(slide, [
        f"目标: {target}",
        f"原型: {proto}",
        f"EEG: {eeg}",
    ], left + Inches(0.15), Inches(2.0), Inches(2.7), 13, LIGHT_GRAY)

# 可调参数
add_shape_bg(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5), BG_LIGHT)
add_text_box(slide, Inches(1.0), Inches(4.7), Inches(5), Inches(0.4), "游戏可调参数空间", 18, ACCENT_BLUE, True)

params = [
    ("难度维度", "项目数量 · 时间压力 · 干扰物 · 相似度 · 序列长度"),
    ("感官维度", "视觉对比度 · 音量 · 动画速度 · 颜色饱和度 · 字体大小"),
    ("反馈维度", "奖励频率 · 奖励类型 · 提示等级 · 错误反馈强度"),
    ("时间维度", "单次时长(5-20min) · 休息间隔 · 总训练时长"),
]

for i, (dim, detail) in enumerate(params):
    y = Inches(5.2) + Inches(0.42) * i
    add_text_box(slide, Inches(1.2), y, Inches(1.8), Inches(0.4), dim, 14, ACCENT_GREEN, True)
    add_text_box(slide, Inches(3.2), y, Inches(9.0), Inches(0.4), detail, 13, LIGHT_GRAY, False)


# ============ Slide 13: 闭环架构 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "06 闭环控制架构", 32, ACCENT_GREEN, True)

# 架构图（文字版）
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(2.8), BG_LIGHT)
add_text_box(slide, Inches(1.0), Inches(1.4), Inches(11.3), Inches(0.4), "系统数据流", 18, ACCENT_BLUE, True)

arch_text = (
    "EEG 硬件 ──(BrainFlow/LSL)──▶ Python 处理端 ──(LSL/UDP)──▶ Unity 游戏\n"
    "                                    │                              │\n"
    "                              特征提取 (2s窗口)               参数更新\n"
    "                              状态估计                        行为数据\n"
    "                              控制决策                           │\n"
    "                                    │                     ◀──────┘\n"
    "                              数据存储 (EDF+ / SQLite)"
)
add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10.5), Inches(1.8), arch_text, 13, ACCENT_BLUE, False, PP_ALIGN.LEFT, "Courier New")

# 三层控制
add_text_box(slide, Inches(0.8), Inches(4.3), Inches(11.3), Inches(0.4), "三层控制策略", 18, ACCENT_GREEN, True)

layers = [
    ("第一层：实时安全守护", "0.5s", "Delta 过高→暂停; 参与度过低→注意力唤醒; 疲劳→休息引导", ACCENT_RED),
    ("第二层：自适应难度调节", "5-10s", "SPR 为主控变量，PID 控制维持 Flow Zone\n目标：TAR/SPR 在个体校准的\"最优挑战区间\"内", ACCENT_BLUE),
    ("第三层：跨会话个性化", "天/周级", "追踪长期进步趋势，调整训练计划\n检测平台期→切换游戏；APF 回升→提升目标区间", ACCENT_GREEN),
]

for i, (name, freq, desc, color) in enumerate(layers):
    y = Inches(4.8) + Inches(0.85) * i
    add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(0.78), BG_LIGHT)
    add_text_box(slide, Inches(1.0), y + Pt(4), Inches(3.0), Inches(0.7), name, 15, color, True)
    add_text_box(slide, Inches(4.0), y + Pt(4), Inches(1.0), Inches(0.7), freq, 13, GRAY, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(5.2), y + Pt(2), Inches(7.0), Inches(0.7), desc, 12, LIGHT_GRAY, False)


# ============ Slide 14: 推荐主控指标 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "闭环系统推荐指标组合", 28, ACCENT_GREEN, True)

# 三列
categories = [
    ("⚡ 实时闭环 (每2s)", ACCENT_BLUE, [
        "1. SPR = (α+β)/(δ+θ)",
        "   综合认知效率，主控难度",
        "",
        "2. APF (Alpha 峰频率)",
        "   个体化校准 + 长期追踪",
        "",
        "3. Delta 功率",
        "   安全守护(疲劳/嗜睡)",
    ]),
    ("🎮 任务嵌入 (游戏内)", ACCENT_GREEN, [
        "4. P300 振幅/潜伏期",
        "   目标检测任务表现",
        "   通关/升级条件",
        "",
        "5. 行为指标",
        "   反应时间",
        "   准确率",
        "   操作模式多样性",
    ]),
    ("📊 阶段评估 (会话前后)", ACCENT_ORANGE, [
        "6. Beta 功率",
        "   参与度/认知活跃度",
        "",
        "7. α/β AEC (连接性)",
        "   网络完整性",
        "   可能是最早的改善信号",
        "",
        "8. 微状态 C",
        "   辅助疗效评估",
    ]),
]

for i, (title, color, items) in enumerate(categories):
    left = Inches(0.6) + Inches(4.2) * i
    add_shape_bg(slide, left, Inches(1.3), Inches(3.9), Inches(5.8), BG_LIGHT)
    add_text_box(slide, left + Inches(0.15), Inches(1.4), Inches(3.6), Inches(0.5), title, 17, color, True, PP_ALIGN.CENTER)
    add_bullet_slide_content(slide, items, left + Inches(0.15), Inches(2.0), Inches(3.6), 14, LIGHT_GRAY)


# ============ Slide 15: 技术栈 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "推荐技术栈", 28, ACCENT_BLUE, True)

tech = [
    ("游戏引擎", "Unity (C#)", "跨平台、生态成熟、EEG SDK 对接方便", ACCENT_BLUE),
    ("EEG 设备", "OpenBCI Cyton 8ch", "开源、可定制、研究级信号质量", ACCENT_GREEN),
    ("数据流", "Lab Streaming Layer (LSL)", "神经科学标准，时间同步精确", ACCENT_ORANGE),
    ("信号处理", "Python (BrainFlow + MNE)", "跨设备统一 API + 成熟特征提取", ACCENT_PURPLE),
    ("控制策略", "PID → XGBoost", "先规则积累数据，后 ML 优化", ACCENT_BLUE),
    ("数据存储", "EDF+ (原始) + SQLite", "标准格式 + 轻量分析数据库", ACCENT_GREEN),
]

for i, (category, tool, reason, color) in enumerate(tech):
    y = Inches(1.3) + Inches(0.9) * i
    add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(0.78), BG_LIGHT)
    add_text_box(slide, Inches(1.0), y + Pt(6), Inches(2.0), Inches(0.6), category, 16, color, True)
    add_text_box(slide, Inches(3.2), y + Pt(6), Inches(3.5), Inches(0.6), tool, 16, WHITE, True)
    add_text_box(slide, Inches(7.0), y + Pt(6), Inches(5.2), Inches(0.6), reason, 14, GRAY, False)


# ============ Slide 16: 下一步 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7), "07 下一步计划", 32, ACCENT_ORANGE, True)

steps = [
    ("近期 (1-2周)", ACCENT_BLUE, [
        "补充 AD/MCI 游戏干预的专项文献调研",
        "确定目标人群（MCI? 轻度AD? 中度AD?）",
        "选定 EEG 设备（OpenBCI / Emotiv / NeuroSky）",
        "搭建 Unity + Python + LSL 基础通信框架",
    ]),
    ("中期 (1-2月)", ACCENT_GREEN, [
        "原型开发：先做 1 个子游戏（建议：记忆花园）",
        "实现 SPR + APF + Delta 的实时提取 Pipeline",
        "实现 PID 自适应难度控制器",
        "健康人测试验证闭环功能",
    ]),
    ("远期 (3-6月)", ACCENT_ORANGE, [
        "MCI/AD 患者初步测试（需伦理审查）",
        "完善多游戏模块",
        "积累数据后切 ML 难度预测",
        "论文撰写（框架已搭好）",
    ]),
]

for i, (phase, color, items) in enumerate(steps):
    left = Inches(0.6) + Inches(4.2) * i
    add_shape_bg(slide, left, Inches(1.3), Inches(3.9), Inches(5.5), BG_LIGHT)
    add_text_box(slide, left + Inches(0.15), Inches(1.4), Inches(3.6), Inches(0.5), phase, 20, color, True, PP_ALIGN.CENTER)
    add_bullet_slide_content(slide, [f"• {x}" for x in items], left + Inches(0.15), Inches(2.1), Inches(3.6), 15, LIGHT_GRAY, Pt(12))


# ============ Slide 17: 致谢/结尾 ============
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(13.333), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
             "谢谢", 48, ACCENT_BLUE, True, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.8),
             "面向 AD 患者的 EEG 闭环游戏干预系统 · 调研汇报", 20, GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.6),
             "核心文献: Babiloni et al. 2021 (专家共识) + Sohrabpour et al. 2025 (系统综述) + 7篇闭环游戏文献",
             14, LIGHT_GRAY, False, PP_ALIGN.CENTER)


# ============ 保存 ============
output_path = "/mnt/eason/Jacob-paper-reading/EEG-Closed-Loop-Game/AD-BCI-Game-Survey/AD脑控游戏调研汇报.pptx"
prs.save(output_path)
print(f"✅ PPT 已保存: {output_path}")
print(f"共 {len(prs.slides)} 页")
